# %%
from os import replace
import time
from selenium.webdriver.chrome.options import Options
from lxml import etree
from bs4 import BeautifulSoup
import requests
import pypatent
from selenium import webdriver
import pptx
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
# %%
# Requires geckodriver in your PATH
driver = webdriver.Chrome(r'C:\Users\JACK\Documents\chromedriver.exe')

conn = pypatent.WebConnection(use_selenium=True, selenium_driver=driver)

q = 'AN/"Infineon Technologies"'

res = pypatent.Search(q.replace(" ", "%20"), results_limit=1,
                      get_patent_details=True, web_connection=conn).as_list()
p_list = []
print(res)
# %%


def get_info(num):
    pnum = "US "+str(num)
    chrome = webdriver.Chrome(r'C:\Users\JACK\Documents\chromedriver.exe')
    # chrome.get("https://patents.google.com/patent/US"+num)
    chrome.get("https://www.patentguru.com/search?q=US"+num)

    soup = BeautifulSoup(chrome.page_source, 'html.parser')

    e = soup.find_all("span", {"class": "relative"})
    x = str(e[0]).split("</span>")
    t = x[1]
    a = t.replace("\n", "")
    pnum = a.replace(" ", "")
    p_list.append(pnum)
    if pnum[0] != "U":
        pnum = "US "+str(num)
    chrome.close()
    return pnum

# %%


def dowload_image(pnum):
    chrome = webdriver.Chrome(r'C:\Users\JACK\Documents\chromedriver.exe')
    chrome.get("https://www.patentguru.com/"+pnum)
    soup = BeautifulSoup(chrome.page_source, 'html.parser')
    p = soup.find_all('img')
    p_link = p[3]["data-original"]
    print(p_link)
    r = requests.get(p_link, allow_redirects=True)
    open(pnum+".png", 'wb').write(r.content)
    chrome.close()


# %%%
pptFile = pptx.Presentation()


def new_info(pnum="", ttl="", fdate="", pdate="", name="", an="", abst=""):
    slide = pptFile.slides.add_slide(pptFile.slide_layouts[6])
    table = slide.shapes.add_table(rows=7, cols=4, left=Inches(
        0.5), top=Inches(0.25), width=Inches(9), height=Inches(6))
    x = table.table
    a = x.cell(1, 0)
    b = x.cell(3, 0)
    a.merge(b)
    a = x.cell(1, 1)
    b = x.cell(3, 1)
    a.merge(b)
    a = x.cell(1, 2)
    b = x.cell(2, 2)
    a.merge(b)
    a = x.cell(1, 3)
    b = x.cell(2, 3)
    a.merge(b)
    a = x.cell(4, 0)
    b = x.cell(6, 3)
    a.merge(b)

    cell = x.cell(0, 0)
    cell.text = "Patent No."
    cell = x.cell(1, 0)
    cell.text = "Title"
    cell = x.cell(0, 2)
    cell.text = "Date of Patent"
    cell = x.cell(1, 2)
    cell.text = "Inventors"
    cell = x.cell(3, 2)
    cell.text = "Assignee"
    cell = x.cell(4, 0)
    cell.text = "Abstract:"

    cell = x.cell(0, 1)
    cell.text = pnum
    cell = x.cell(1, 1)
    cell.text = ttl
    cell = x.cell(0, 3)
    cell.text = "Filed: " + fdate+"\n"+"Published: "+pdate
    cell = x.cell(1, 3)
    cell.text = name
    cell = x.cell(3, 3)
    cell.text = an
    cell = x.cell(4, 0)
    cell.text = "Abstract: \n"+abst

    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(x):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)


def new_claims(claims=""):
    # set slide layout
    slide = pptFile.slides.add_slide(pptFile.slide_layouts[5])
    # crete new page
    title = slide.shapes.title
    title.text = "Claims"
    title.top = Inches(0.5)
    title.width = Inches(5)
    title.left = Inches(2.3)
    left, top, width, height = Cm(1), Cm(2), Cm(24), Cm(50)
    # 添加文字段落
    text_frame = slide.shapes.add_textbox(
        left=left, top=top, width=width, height=height).text_frame
    text_frame.text = claims
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(13)
    text_frame.margin_bottom = Inches(0.08)
    text_frame.margin_left = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


# %%
with open('patent.txt', 'w') as f:
    for n in range(len(res)):
        f.write(res[n]["title"]+"\r\n")
        ttl = res[n]["title"]
        num = res[n]["patent_num"].replace(',', "")
        pnum = get_info(num)
        dowload_image(pnum)
        print(pnum)
        print(ttl)
        f.write(res[n]["assignee_name"]+"\r\n")
        an = res[n]["assignee_name"]+", ("+res[n]["assignee_loc"]+")"
        name = ""
        for i in res[n]["inventors"]:
            count = 0
            for t in i:
                if count == 0:
                    f.write(str(t)+", ")
                    name += (str(t)+", ")
                elif count == 1:
                    f.write(str(t)+" (")
                    name += (str(t)+" (")
                elif count == 2:
                    f.write(str(t)+")")
                    name += (str(t)+")")
                count += 1
            name += "\n"
            f.write("\r\n")
        f.write(res[n]["patent_date"]+"\r\n")
        pdate = res[n]["patent_date"]
        f.write(res[n]["file_date"]+"\r\n")
        fdate = res[n]["file_date"]
        f.write(res[n]["abstract"]+"\r\n")
        abst = res[n]["abstract"]
        claims = ""
        for c in res[n]["claims"]:
            claims += (c+"\n")
            f.write(str(c)+"\r\n")
        new_info(pnum, ttl, fdate, pdate, name, an, abst)
        new_claims(claims)
    pptFile.save('test.pptx')


# %%
def dowload_image(pnum):
    chrome = webdriver.Chrome(r'C:\Users\JACK\Documents\chromedriver.exe')
    chrome.get("https://www.patentguru.com/"+pnum)
    soup = BeautifulSoup(chrome.page_source, 'html.parser')
    p = soup.find_all('img')
    p_link = (p[3])["data-original"]
    print(p_link)
    r = requests.get(p_link, allow_redirects=True)
    open(pnum+".png", 'wb').write(r.content)
    chrome.close()


p_list = ['US10404256B2']
print(p_list)
for p in p_list:
    dowload_image(p)

# %%
